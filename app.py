from flask import Flask, render_template, request, send_file
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
import io

app = Flask(__name__)

def clean_text(text):
    """Remove unwanted characters from the text."""
    return text.replace('_x000D_', '').strip()

@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        title = request.form.get("title", "")
        content = request.form.get("content", "")
        lines_per_slide = int(request.form.get("lines_per_slide", 5))
        font_size = int(request.form.get("font_size", 14))
        title_font_size = int(request.form.get("title_font_size", 96))

        # Caminhos para as imagens
        background_image_path = 'static/img/background.jpg'
        logo_image_path = 'static/img/logo.png'

        # Clean the content
        content = clean_text(content)

        # Convert title and content to uppercase
        title = title.upper()
        content = content.upper()

        # Split content into lines
        content_lines = content.splitlines()

        # Create a PowerPoint presentation
        prs = Presentation()
        prs.slide_width = Inches(13.33)  # Increase slide width
        prs.slide_height = Inches(7.5)   # Keep default height

        def set_slide_background(slide, image_path):
            """Set the background image for a slide."""
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            slide.shapes.add_picture(image_path, 0, 0, slide_width, slide_height)

        def add_logo_and_text(slide, logo_path):
            """Add logo and text at the bottom left of the slide."""
            logo_left = Inches(0.2)
            logo_top = prs.slide_height - Inches(1.5)
            logo_width = Inches(1.2)
            slide.shapes.add_picture(logo_path, logo_left, logo_top, width=logo_width)

            text_left = logo_left
            text_top = logo_top + Inches(1.0)
            text_width = logo_width
            text_height = Inches(0.5)
            textbox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
            text_frame = textbox.text_frame
            text_frame.text = "PIBO"

            p = text_frame.paragraphs[0]
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER

        # Title slide
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, background_image_path)
        add_logo_and_text(slide, logo_image_path)
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        title_shape = slide.shapes.add_textbox(Inches(0), Inches(0), slide_width, slide_height)
        text_frame = title_shape.text_frame
        text_frame.text = title

        title_paragraph = text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(title_font_size)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(255, 255, 255)
        title_paragraph.font.name = "Calibri"
        title_paragraph.alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        # Add content slides
        for i in range(0, len(content_lines), lines_per_slide):
            content_slide = prs.slides.add_slide(prs.slide_layouts[5])
            set_slide_background(content_slide, background_image_path)
            add_logo_and_text(content_slide, logo_image_path)

            content_shape = content_slide.shapes.add_textbox(Inches(0), Inches(0), slide_width, slide_height)
            text_frame = content_shape.text_frame
            text_frame.text = "\n".join(content_lines[i:i + lines_per_slide])

            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.font.name = "Calibri"
                paragraph.font.bold = True
                paragraph.alignment = PP_ALIGN.CENTER

            text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        # Save presentation to a BytesIO object
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)

        return send_file(output, as_attachment=True, download_name="presentation.pptx", mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    return render_template("index.html")

if __name__ == "__main__":
    app.run(debug=True)
